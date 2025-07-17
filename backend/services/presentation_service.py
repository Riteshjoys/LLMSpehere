import asyncio
import json
import base64
from typing import List, Dict, Any, Optional
from datetime import datetime
from io import BytesIO
import uuid
import os
import io

# Presentation libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.lineplots import LinePlot
from reportlab.graphics.charts.barcharts import VerticalBarChart

# Data visualization
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np

# Google Slides API (placeholder for future implementation)
try:
    from google.oauth2.service_account import Credentials
    from googleapiclient.discovery import build
    GOOGLE_SLIDES_AVAILABLE = True
except ImportError:
    GOOGLE_SLIDES_AVAILABLE = False

from models.presentation_models import (
    PresentationRequest, PresentationResponse, PresentationGeneration,
    SlideContent, SlideType, ChartType, PresentationFormat, PresentationTheme,
    SlideGenerationRequest, ChartData, PresentationTemplate, SlideLayout
)
from services.text_generation_service import TextGenerationService
from services.image_generation_service import ImageGenerationService
from utils.database import get_database

class PresentationService:
    def __init__(self):
        self.db = get_database()
        self.text_service = TextGenerationService()
        self.image_service = ImageGenerationService()
        self.templates_collection = "presentation_templates"
        self.presentations_collection = "presentations"
        self.slides_collection = "presentation_slides"
        self.history_collection = "presentation_history"
        
        # Default templates
        self.default_templates = [
            {
                "id": "business_pitch",
                "name": "Business Pitch Deck",
                "description": "Professional pitch deck template for business presentations",
                "type": "business",
                "slides": [
                    {"title": "Title Slide", "layout": "title"},
                    {"title": "Problem Statement", "layout": "content"},
                    {"title": "Solution", "layout": "content"},
                    {"title": "Market Analysis", "layout": "chart"},
                    {"title": "Business Model", "layout": "content"},
                    {"title": "Financial Projections", "layout": "chart"},
                    {"title": "Team", "layout": "content"},
                    {"title": "Ask & Next Steps", "layout": "content"}
                ]
            },
            {
                "id": "marketing_report",
                "name": "Marketing Report",
                "description": "Template for marketing performance reports",
                "type": "marketing",
                "slides": [
                    {"title": "Executive Summary", "layout": "title"},
                    {"title": "Key Metrics", "layout": "chart"},
                    {"title": "Campaign Performance", "layout": "chart"},
                    {"title": "Audience Analysis", "layout": "content"},
                    {"title": "ROI Analysis", "layout": "chart"},
                    {"title": "Recommendations", "layout": "content"},
                    {"title": "Next Steps", "layout": "content"}
                ]
            },
            {
                "id": "product_demo",
                "name": "Product Demo",
                "description": "Template for product demonstration presentations",
                "type": "product",
                "slides": [
                    {"title": "Product Overview", "layout": "title"},
                    {"title": "Key Features", "layout": "content"},
                    {"title": "Benefits", "layout": "content"},
                    {"title": "Use Cases", "layout": "content"},
                    {"title": "Demo", "layout": "content"},
                    {"title": "Pricing", "layout": "content"},
                    {"title": "Q&A", "layout": "content"}
                ]
            }
        ]

    async def generate_presentation(self, request: PresentationRequest, user_id: str) -> PresentationResponse:
        """Generate a complete presentation"""
        try:
            # Generate presentation outline
            outline = await self._generate_outline(request)
            
            # Generate slides based on outline
            slides = await self._generate_slides(outline, request, user_id)
            
            # Create presentation record
            presentation = PresentationGeneration(
                presentation_id=str(uuid.uuid4()),
                user_id=user_id,
                title=request.title,
                topic=request.topic,
                slides=slides,
                theme=request.theme,
                format=request.format,
                total_slides=len(slides),
                provider=request.provider_name,
                model=request.model
            )
            
            # Save to database
            await self._save_presentation(presentation)
            
            # Generate file if requested
            file_data = await self._generate_file(presentation)
            
            return PresentationResponse(
                presentation_id=presentation.presentation_id,
                title=presentation.title,
                topic=presentation.topic,
                slides=presentation.slides,
                theme=presentation.theme,
                format=presentation.format,
                total_slides=presentation.total_slides,
                file_base64=file_data.get('base64') if file_data else None,
                file_url=file_data.get('url') if file_data else None,
                thumbnail_url=file_data.get('thumbnail') if file_data else None,
                created_at=presentation.created_at,
                user_id=presentation.user_id,
                provider=presentation.provider,
                model=presentation.model
            )
            
        except Exception as e:
            raise Exception(f"Failed to generate presentation: {str(e)}")

    async def get_templates(self, db=None):
        """Get all available templates"""
        try:
            if db is None:
                db = self.db
            
            # Get custom templates from database
            custom_templates = []
            try:
                for template in db[self.templates_collection].find():
                    template['_id'] = str(template['_id'])
                    custom_templates.append(template)
            except:
                pass
            
            # Combine with default templates
            all_templates = self.default_templates + custom_templates
            return all_templates
        except Exception as e:
            raise Exception(f"Error getting templates: {str(e)}")

    async def create_template(self, db, name: str, description: str, template_type: str, 
                            file, user_id: str):
        """Create a new presentation template"""
        try:
            # Read uploaded file
            file_content = await file.read()
            
            # Create template document
            template_doc = {
                "id": str(uuid.uuid4()),
                "name": name,
                "description": description,
                "type": template_type,
                "file_content": base64.b64encode(file_content).decode('utf-8'),
                "file_name": file.filename,
                "file_size": len(file_content),
                "created_by": user_id,
                "created_at": datetime.utcnow(),
                "updated_at": datetime.utcnow()
            }
            
            # Insert into database
            result = db[self.templates_collection].insert_one(template_doc)
            return template_doc["id"]
        except Exception as e:
            raise Exception(f"Error creating template: {str(e)}")

    async def get_template(self, db, template_id: str):
        """Get a specific template by ID"""
        try:
            # Check default templates first
            for template in self.default_templates:
                if template["id"] == template_id:
                    return template
            
            # Check custom templates
            template = db[self.templates_collection].find_one({"id": template_id})
            if template:
                template['_id'] = str(template['_id'])
                return template
            
            return None
        except Exception as e:
            raise Exception(f"Error getting template: {str(e)}")

    async def create_presentation(self, db, template_id: str, title: str, 
                                data: Dict[str, Any], user_id: str):
        """Create a new presentation from template"""
        try:
            # Get template
            template = await self.get_template(db, template_id)
            if not template:
                raise Exception("Template not found")
            
            # Create presentation document
            presentation_doc = {
                "id": str(uuid.uuid4()),
                "title": title,
                "template_id": template_id,
                "user_id": user_id,
                "data": data,
                "slides": template.get("slides", []),
                "created_at": datetime.utcnow(),
                "updated_at": datetime.utcnow(),
                "status": "draft"
            }
            
            # Insert into database
            result = db[self.presentations_collection].insert_one(presentation_doc)
            
            # Save to history
            await self._save_to_history(db, presentation_doc["id"], user_id, "created")
            
            return presentation_doc["id"]
        except Exception as e:
            raise Exception(f"Error creating presentation: {str(e)}")

    async def get_user_presentations(self, db, user_id: str):
        """Get all presentations for a user"""
        try:
            presentations = []
            for presentation in db[self.presentations_collection].find({"user_id": user_id}):
                presentation['_id'] = str(presentation['_id'])
                presentations.append(presentation)
            return presentations
        except Exception as e:
            raise Exception(f"Error getting user presentations: {str(e)}")

    async def get_presentation_by_id(self, db, presentation_id: str):
        """Get a specific presentation by ID (database version)"""
        try:
            presentation = db[self.presentations_collection].find_one({"id": presentation_id})
            if presentation:
                presentation['_id'] = str(presentation['_id'])
            return presentation
        except Exception as e:
            raise Exception(f"Error getting presentation: {str(e)}")

    async def update_presentation(self, db, presentation_id: str, updates: Dict[str, Any]):
        """Update a presentation"""
        try:
            updates["updated_at"] = datetime.utcnow()
            result = db[self.presentations_collection].update_one(
                {"id": presentation_id}, 
                {"$set": updates}
            )
            return result.modified_count > 0
        except Exception as e:
            raise Exception(f"Error updating presentation: {str(e)}")

    async def export_presentation(self, db, presentation_id: str, format: str):
        """Export presentation in specified format"""
        try:
            # Get presentation
            presentation = await self.get_presentation_by_id(db, presentation_id)
            if not presentation:
                raise Exception("Presentation not found")
            
            if format == "pptx":
                return await self._export_to_pptx(presentation)
            elif format == "pdf":
                return await self._export_to_pdf(presentation)
            elif format == "google-slides":
                return await self._export_to_google_slides(presentation)
            else:
                raise Exception("Unsupported export format")
        except Exception as e:
            raise Exception(f"Error exporting presentation: {str(e)}")

    async def generate_presentation_content(self, db, presentation_id: str, request: Dict[str, Any]):
        """Generate presentation content using AI"""
        try:
            # Get presentation
            presentation = await self.get_presentation_by_id(db, presentation_id)
            if not presentation:
                raise Exception("Presentation not found")
            
            # Generate content based on request
            content_type = request.get("content_type", "text")
            prompt = request.get("prompt", "")
            
            if content_type == "text":
                # Generate text content
                generated_content = await self._generate_text_content(prompt)
            elif content_type == "chart":
                # Generate chart data
                generated_content = await self._generate_chart_data(request)
            else:
                raise Exception("Unsupported content type")
            
            return {"content": generated_content, "type": content_type}
        except Exception as e:
            raise Exception(f"Error generating content: {str(e)}")

    async def delete_presentation(self, db, presentation_id: str):
        """Delete a presentation"""
        try:
            result = db[self.presentations_collection].delete_one({"id": presentation_id})
            return result.deleted_count > 0
        except Exception as e:
            raise Exception(f"Error deleting presentation: {str(e)}")

    async def add_slide(self, db, presentation_id: str, slide_data: Dict[str, Any]):
        """Add a new slide to presentation"""
        try:
            slide_id = str(uuid.uuid4())
            slide_data["id"] = slide_id
            
            # Add slide to presentation
            result = db[self.presentations_collection].update_one(
                {"id": presentation_id},
                {"$push": {"slides": slide_data}, "$set": {"updated_at": datetime.utcnow()}}
            )
            
            return slide_id
        except Exception as e:
            raise Exception(f"Error adding slide: {str(e)}")

    async def update_slide(self, db, presentation_id: str, slide_id: str, updates: Dict[str, Any]):
        """Update a slide in presentation"""
        try:
            result = db[self.presentations_collection].update_one(
                {"id": presentation_id, "slides.id": slide_id},
                {"$set": {f"slides.$.{k}": v for k, v in updates.items()}}
            )
            
            db[self.presentations_collection].update_one(
                {"id": presentation_id},
                {"$set": {"updated_at": datetime.utcnow()}}
            )
            
            return result.modified_count > 0
        except Exception as e:
            raise Exception(f"Error updating slide: {str(e)}")

    async def delete_slide(self, db, presentation_id: str, slide_id: str):
        """Delete a slide from presentation"""
        try:
            result = db[self.presentations_collection].update_one(
                {"id": presentation_id},
                {"$pull": {"slides": {"id": slide_id}}, "$set": {"updated_at": datetime.utcnow()}}
            )
            
            return result.modified_count > 0
        except Exception as e:
            raise Exception(f"Error deleting slide: {str(e)}")

    async def create_chart(self, db, presentation_id: str, chart_request: Dict[str, Any]):
        """Create a chart for presentation"""
        try:
            chart_type = chart_request.get("type", "bar")
            chart_data = chart_request.get("data", {})
            
            # Generate chart image
            chart_image = await self._generate_chart_image(chart_type, chart_data)
            
            return {
                "chart_type": chart_type,
                "chart_data": chart_data,
                "chart_image": chart_image
            }
        except Exception as e:
            raise Exception(f"Error creating chart: {str(e)}")

    async def get_presentation_history(self, db, user_id: str):
        """Get presentation history for user"""
        try:
            history = []
            for item in db[self.history_collection].find({"user_id": user_id}).sort("created_at", -1):
                item['_id'] = str(item['_id'])
                history.append(item)
            return history
        except Exception as e:
            raise Exception(f"Error getting presentation history: {str(e)}")

    async def get_presentation_stats(self, db, user_id: str):
        """Get presentation statistics for user"""
        try:
            # Count total presentations
            total_presentations = db[self.presentations_collection].count_documents({"user_id": user_id})
            
            # Count by template type
            pipeline = [
                {"$match": {"user_id": user_id}},
                {"$group": {"_id": "$template_id", "count": {"$sum": 1}}}
            ]
            type_counts = {}
            for result in db[self.presentations_collection].aggregate(pipeline):
                type_counts[result["_id"]] = result["count"]
            
            # Recent activity
            recent_activity = []
            for item in db[self.history_collection].find({"user_id": user_id}).sort("created_at", -1).limit(10):
                item['_id'] = str(item['_id'])
                recent_activity.append(item)
            
            return {
                "total_presentations": total_presentations,
                "type_counts": type_counts,
                "recent_activity": recent_activity
            }
        except Exception as e:
            raise Exception(f"Error getting presentation stats: {str(e)}")

    # Private helper methods from original implementation
    async def _generate_outline(self, request: PresentationRequest) -> List[str]:
        """Generate presentation outline using AI"""
        if request.outline:
            return request.outline
        
        prompt = f"""
        Create a detailed outline for a {request.num_slides}-slide presentation on "{request.topic}".
        
        Title: {request.title}
        Audience: {request.audience}
        Tone: {request.tone}
        
        Please provide a JSON array of slide titles/topics that would make an effective presentation.
        Each slide should be a string describing the main point or topic for that slide.
        
        The first slide should be a title slide, and the last slide should be a conclusion or thank you slide.
        
        Example format:
        [
            "Title: {request.title}",
            "Introduction to {request.topic}",
            "Key Benefits",
            "Implementation Strategy",
            "Case Studies",
            "Next Steps",
            "Thank You & Questions"
        ]
        """
        
        try:
            # Use text generation service to create outline
            outline_response = await self.text_service.generate_text(
                provider_name=request.provider_name or "openai",
                model=request.model or "gpt-4",
                prompt=prompt,
                max_tokens=1000,
                user_id=request.user_id if hasattr(request, 'user_id') else "system"
            )
            
            # Parse the response to extract outline
            outline_text = outline_response.get('content', '').strip()
            
            # Try to parse as JSON first
            try:
                outline = json.loads(outline_text)
                if isinstance(outline, list):
                    return outline[:request.num_slides]
            except:
                # If not JSON, split by lines and clean up
                lines = outline_text.split('\n')
                outline = []
                for line in lines:
                    line = line.strip()
                    if line and not line.startswith(('```', 'json', 'Here')):
                        # Clean up line formatting
                        line = line.replace('"', '').replace(',', '').strip()
                        if line.startswith(('- ', '* ', '1. ', '2. ', '3. ', '4. ', '5. ', '6. ', '7. ', '8. ', '9. ')):
                            line = line[2:].strip()
                        outline.append(line)
                
                return outline[:request.num_slides]
            
        except Exception as e:
            # Fallback to default outline
            return self._get_default_outline(request)

    def _get_default_outline(self, request: PresentationRequest) -> List[str]:
        """Generate a default outline if AI generation fails"""
        base_outline = [
            f"Title: {request.title}",
            f"Introduction to {request.topic}",
            "Key Points",
            "Benefits and Advantages",
            "Implementation",
            "Examples and Case Studies",
            "Challenges and Solutions",
            "Future Outlook",
            "Conclusion",
            "Thank You & Questions"
        ]
        
        return base_outline[:request.num_slides]

    async def _generate_slides(self, outline: List[str], request: PresentationRequest, user_id: str) -> List[SlideContent]:
        """Generate slide content for each outline item"""
        slides = []
        
        for i, slide_title in enumerate(outline):
            slide_type = self._determine_slide_type(slide_title, i, len(outline))
            
            slide_content = await self._generate_slide_content(
                slide_title=slide_title,
                slide_type=slide_type,
                request=request,
                user_id=user_id,
                position=i
            )
            
            slides.append(slide_content)
        
        return slides

    def _determine_slide_type(self, title: str, position: int, total_slides: int) -> SlideType:
        """Determine slide type based on title and position"""
        title_lower = title.lower()
        
        if position == 0 or "title" in title_lower:
            return SlideType.TITLE
        elif position == total_slides - 1 or any(word in title_lower for word in ["conclusion", "thank", "questions", "summary"]):
            return SlideType.CONCLUSION
        elif any(word in title_lower for word in ["chart", "data", "statistics", "metrics", "numbers"]):
            return SlideType.CHART
        elif any(word in title_lower for word in ["image", "example", "case", "demo"]):
            return SlideType.IMAGE
        elif any(word in title_lower for word in ["vs", "comparison", "compare", "differences"]):
            return SlideType.COMPARISON
        else:
            return SlideType.CONTENT

    async def _generate_slide_content(self, slide_title: str, slide_type: SlideType, request: PresentationRequest, user_id: str, position: int) -> SlideContent:
        """Generate content for a specific slide"""
        
        # Create base slide
        slide = SlideContent(
            type=slide_type,
            layout_id=self._get_layout_for_type(slide_type),
            title=slide_title,
            position=position
        )
        
        if slide_type == SlideType.TITLE:
            slide.title = request.title
            slide.content = f"Presentation on {request.topic}"
            slide.speaker_notes = f"Welcome to this presentation on {request.topic}. Today we'll explore the key aspects and insights."
        
        elif slide_type == SlideType.CONCLUSION:
            slide.title = "Conclusion"
            slide.content = await self._generate_conclusion_content(request, user_id)
            slide.speaker_notes = "Thank you for your attention. Are there any questions?"
        
        else:
            # Generate detailed content for regular slides
            content_prompt = f"""
            Create content for a presentation slide with the title: "{slide_title}"
            
            Context:
            - Overall presentation topic: {request.topic}
            - Audience: {request.audience}
            - Tone: {request.tone}
            
            Please provide:
            1. Main content/bullet points (3-5 key points)
            2. Speaker notes (2-3 sentences for the presenter)
            
            Format the response as:
            CONTENT:
            • Point 1
            • Point 2
            • Point 3
            
            SPEAKER_NOTES:
            Notes for the presenter...
            """
            
            try:
                content_response = await self.text_service.generate_text(
                    provider_name=request.provider_name or "openai",
                    model=request.model or "gpt-4",
                    prompt=content_prompt,
                    max_tokens=500,
                    user_id=user_id
                )
                
                content_text = content_response.get('content', '').strip()
                
                # Parse content and speaker notes
                if "CONTENT:" in content_text and "SPEAKER_NOTES:" in content_text:
                    parts = content_text.split("SPEAKER_NOTES:")
                    content_part = parts[0].replace("CONTENT:", "").strip()
                    notes_part = parts[1].strip()
                    
                    # Extract bullet points
                    bullet_points = []
                    for line in content_part.split('\n'):
                        line = line.strip()
                        if line and line.startswith(('•', '-', '*')):
                            bullet_points.append(line[1:].strip())
                    
                    slide.bullet_points = bullet_points
                    slide.content = content_part
                    slide.speaker_notes = notes_part
                else:
                    slide.content = content_text
                    slide.speaker_notes = f"Present the content about {slide_title}."
                
            except Exception as e:
                # Fallback content
                slide.content = f"Key points about {slide_title} will be discussed here."
                slide.speaker_notes = f"Discuss the main aspects of {slide_title}."
        
        # Generate image if requested and appropriate
        if request.include_images and slide_type in [SlideType.CONTENT, SlideType.IMAGE]:
            try:
                image_data = await self._generate_slide_image(slide_title, request, user_id)
                if image_data:
                    slide.image_base64 = image_data
            except Exception as e:
                print(f"Failed to generate image for slide: {e}")
        
        # Generate chart if requested and appropriate
        if request.include_charts and slide_type == SlideType.CHART:
            try:
                chart_data = await self._generate_chart_data(slide_title, request)
                if chart_data:
                    slide.chart_data = chart_data
                    slide.chart_type = ChartType.BAR  # Default
            except Exception as e:
                print(f"Failed to generate chart for slide: {e}")
        
        return slide

    async def _generate_conclusion_content(self, request: PresentationRequest, user_id: str) -> str:
        """Generate conclusion content"""
        prompt = f"""
        Create a conclusion for a presentation on "{request.topic}".
        
        Title: {request.title}
        Audience: {request.audience}
        
        Please provide a brief conclusion that summarizes the key takeaways and includes a call to action.
        Keep it concise and impactful.
        """
        
        try:
            response = await self.text_service.generate_text(
                provider_name=request.provider_name or "openai",
                model=request.model or "gpt-4",
                prompt=prompt,
                max_tokens=200,
                user_id=user_id
            )
            return response.get('content', 'Thank you for your attention. Questions?')
        except:
            return 'Thank you for your attention. Questions?'

    async def _generate_slide_image(self, slide_title: str, request: PresentationRequest, user_id: str) -> Optional[str]:
        """Generate an image for a slide"""
        image_prompt = f"Professional illustration for presentation slide about {slide_title}, {request.tone} style, business presentation format"
        
        try:
            image_response = await self.image_service.generate_image(
                provider_name="openai",  # Default to OpenAI for now
                model="dall-e-3",
                prompt=image_prompt,
                user_id=user_id
            )
            return image_response.get('image_base64')
        except Exception as e:
            print(f"Image generation failed: {e}")
            return None

    async def _generate_chart_data(self, slide_title: str, request: PresentationRequest) -> Optional[Dict[str, Any]]:
        """Generate sample chart data for a slide"""
        # This would integrate with a data visualization service
        # For now, return sample data
        return {
            "labels": ["Q1", "Q2", "Q3", "Q4"],
            "datasets": [{
                "label": "Performance",
                "data": [65, 59, 80, 81],
                "backgroundColor": "#4F46E5"
            }]
        }

    def _get_layout_for_type(self, slide_type: SlideType) -> str:
        """Get appropriate layout ID for slide type"""
        layout_map = {
            SlideType.TITLE: "title_layout",
            SlideType.CONTENT: "content_layout",
            SlideType.IMAGE: "image_layout",
            SlideType.CHART: "chart_layout",
            SlideType.COMPARISON: "comparison_layout",
            SlideType.CONCLUSION: "conclusion_layout"
        }
        return layout_map.get(slide_type, "content_layout")

    async def _save_presentation(self, presentation: PresentationGeneration):
        """Save presentation to database"""
        try:
            presentation_data = presentation.dict()
            presentation_data['_id'] = presentation.presentation_id
            
            await self.db.presentation_generations.insert_one(presentation_data)
        except Exception as e:
            print(f"Failed to save presentation: {e}")

    async def _generate_file(self, presentation: PresentationGeneration) -> Optional[Dict[str, Any]]:
        """Generate presentation file (placeholder for now)"""
        # This would integrate with presentation generation libraries
        # For now, return placeholder data
        return {
            "base64": "placeholder_base64_data",
            "url": f"/presentations/{presentation.presentation_id}.pptx",
            "thumbnail": f"/presentations/{presentation.presentation_id}_thumb.png"
        }

    async def _generate_text_content(self, prompt: str):
        """Generate text content using AI"""
        # Placeholder for AI text generation
        # In a real implementation, this would call an AI service
        return f"Generated content for: {prompt}"

    async def _export_to_pptx(self, presentation: Dict[str, Any]):
        """Export presentation to PowerPoint format"""
        try:
            # Create new presentation
            prs = Presentation()
            
            # Process slides
            for slide_data in presentation.get("slides", []):
                slide_layout = prs.slide_layouts[0]  # Title slide
                if slide_data.get("layout") == "content":
                    slide_layout = prs.slide_layouts[1]  # Content slide
                elif slide_data.get("layout") == "chart":
                    slide_layout = prs.slide_layouts[5]  # Chart slide
                
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")
                
                # Add content
                if len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                    if content_placeholder.has_text_frame:
                        content_placeholder.text = slide_data.get("content", "")
            
            # Save to BytesIO
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            raise Exception(f"Error exporting to PPTX: {str(e)}")

    async def _export_to_pdf(self, presentation: Dict[str, Any]):
        """Export presentation to PDF format"""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            # Title page
            title = Paragraph(presentation.get("title", "Presentation"), styles['Title'])
            elements.append(title)
            elements.append(Spacer(1, 12))
            
            # Process slides
            for slide_data in presentation.get("slides", []):
                # Slide title
                slide_title = Paragraph(slide_data.get("title", ""), styles['Heading1'])
                elements.append(slide_title)
                elements.append(Spacer(1, 12))
                
                # Slide content
                content = slide_data.get("content", "")
                if content:
                    slide_content = Paragraph(content, styles['Normal'])
                    elements.append(slide_content)
                
                elements.append(Spacer(1, 24))
            
            doc.build(elements)
            buffer.seek(0)
            
            return buffer.getvalue()
        except Exception as e:
            raise Exception(f"Error exporting to PDF: {str(e)}")

    async def _export_to_google_slides(self, presentation: Dict[str, Any]):
        """Export presentation to Google Slides"""
        try:
            if not GOOGLE_SLIDES_AVAILABLE:
                raise Exception("Google Slides API not available")
            
            # Placeholder for Google Slides API implementation
            # This would require proper Google API credentials
            return {
                "url": "https://docs.google.com/presentation/d/placeholder",
                "presentation_id": "placeholder"
            }
        except Exception as e:
            raise Exception(f"Error exporting to Google Slides: {str(e)}")

    async def _generate_chart_image(self, chart_type: str, data: Dict[str, Any]):
        """Generate chart image"""
        try:
            # Create matplotlib chart
            fig, ax = plt.subplots(figsize=(10, 6))
            
            categories = data.get("categories", [])
            values = data.get("values", [])
            
            if chart_type == "bar":
                ax.bar(categories, values)
            elif chart_type == "line":
                ax.plot(categories, values, marker='o')
            elif chart_type == "pie":
                ax.pie(values, labels=categories, autopct='%1.1f%%')
            
            ax.set_title(data.get("title", "Chart"))
            
            # Save to base64
            buffer = io.BytesIO()
            plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
            buffer.seek(0)
            
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            plt.close()
            
            return image_base64
        except Exception as e:
            raise Exception(f"Error generating chart: {str(e)}")

    async def _save_to_history(self, db, presentation_id: str, user_id: str, action: str):
        """Save action to history"""
        try:
            history_doc = {
                "id": str(uuid.uuid4()),
                "presentation_id": presentation_id,
                "user_id": user_id,
                "action": action,
                "created_at": datetime.utcnow()
            }
            
            db[self.history_collection].insert_one(history_doc)
        except Exception as e:
            # Don't raise exception for history logging failures
            print(f"Error saving to history: {str(e)}")